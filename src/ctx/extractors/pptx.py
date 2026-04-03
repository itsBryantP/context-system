"""PowerPoint extractor — python-pptx, with LibreOffice noted as a future option."""

from __future__ import annotations

from pathlib import Path

from ctx.extractors.base import Extractor
from ctx.schema import Source, SourceType

_IMPORT_MSG = (
    "python-pptx is required for PPTX extraction. "
    "Install with: uv pip install 'ctx-modules[extractors]'"
)


class PPTXExtractor(Extractor):
    def can_handle(self, source: Source) -> bool:
        return source.type == SourceType.PPTX

    def extract(self, source: Source, output_dir: Path) -> list[Path]:
        if not source.path:
            raise ValueError("PPTX source requires a 'path'")

        pptx_path = Path(source.path).expanduser().resolve()
        if not pptx_path.exists():
            raise FileNotFoundError(f"PPTX not found: {pptx_path}")

        output_dir.mkdir(parents=True, exist_ok=True)
        out_path = output_dir / (pptx_path.stem + ".md")
        out_path.write_text(_extract_pptx(pptx_path))
        return [out_path]


def _extract_pptx(pptx_path: Path) -> str:
    try:
        from pptx import Presentation  # noqa: PLC0415
    except ImportError:
        raise ImportError(_IMPORT_MSG)

    prs = Presentation(str(pptx_path))
    md_lines = [f"# {pptx_path.stem}\n"]

    for i, slide in enumerate(prs.slides, 1):
        md_lines.append(f"\n## Slide {i}")

        title_shape = slide.shapes.title
        if title_shape and title_shape.text.strip():
            md_lines.append(f"\n### {title_shape.text.strip()}")

        for shape in slide.shapes:
            if shape is title_shape or not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if not text:
                    continue
                is_bold = any(
                    run.font.bold
                    for run in para.runs
                    if run.font.bold is not None
                )
                if is_bold:
                    md_lines.append(f"\n**{text}**")
                else:
                    md_lines.append(f"- {text}")

        # Speaker notes as blockquotes
        if slide.has_notes_slide:
            notes_tf = slide.notes_slide.notes_text_frame
            notes_text = notes_tf.text.strip() if notes_tf else ""
            if notes_text:
                md_lines.append("")
                for line in notes_text.splitlines():
                    line = line.strip()
                    if line:
                        md_lines.append(f"> {line}")

    return "\n".join(md_lines) + "\n"
